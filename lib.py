"""
ppt-gen コアライブラリ
汎用PPTXスライド生成ツール — ブランド制約付き
"""
from __future__ import annotations

import io
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# ブランド定数
# ---------------------------------------------------------------------------

FORBIDDEN_FONTS = {"Yu Gothic", "游ゴシック", "MS Gothic", "ＭＳ ゴシック", "MS Mincho", "ＭＳ 明朝"}

C = {
    "text":      "1A1A1A",
    "textLight": "666666",
    "textMuted": "999999",
    "bg":        "FFFFFF",
    "bgLight":   "F5F5F5",
    "accent":    "0070C0",
    "accent2":   "00B050",
    "accent3":   "FF6B00",
    "border":    "D9D9D9",
    "headerBg":  "0070C0",
    "headerText":"FFFFFF",
    "white":     "FFFFFF",
    "black":     "000000",
    "rowAlt":    "EFF4FB",
}
_C_DEFAULT = C.copy()  # テーマリセット用のデフォルト値

# ---------------------------------------------------------------------------
# テーマプリセット
# ---------------------------------------------------------------------------

THEMES: dict[str, dict] = {
    "default": {},  # lib.py のデフォルト値をそのまま使用
    "accenture": {
        "accent":    "A100FF",
        "accent2":   "7300B3",
        "accent3":   "FF6B00",
        "headerBg":  "A100FF",
        "rowAlt":    "F5E6FF",
    },
    "navy": {
        "accent":    "1A3C6E",
        "accent2":   "2E6DB4",
        "accent3":   "C8A951",
        "headerBg":  "1A3C6E",
        "rowAlt":    "EBF0F8",
    },
    "green": {
        "accent":    "00875A",
        "accent2":   "005C3C",
        "accent3":   "F4A100",
        "headerBg":  "00875A",
        "rowAlt":    "E6F5EF",
    },
    "warm": {
        "accent":    "C55A11",
        "accent2":   "843D0B",
        "accent3":   "F4C430",
        "headerBg":  "C55A11",
        "rowAlt":    "FDF0E6",
    },
    "mckinsey": {
        "accent":    "002F6C",
        "accent2":   "0060A9",
        "accent3":   "C8A951",
        "headerBg":  "002F6C",
        "rowAlt":    "E8EDF5",
        "bgLight":   "F2F4F8",
    },
}

TEXT_STYLES: dict[str, dict] = {
    "heading":    {"font_size": 24, "bold": True,  "color": "text"},
    "subheading": {"font_size": 18, "bold": True,  "color": "text"},
    "body":       {"font_size": 14, "bold": False, "color": "text"},
    "small":      {"font_size": 11, "bold": False, "color": "textLight"},
    "caption":    {"font_size": 10, "bold": False, "color": "textMuted"},
    "kpi":        {"font_size": 40, "bold": True,  "color": "text"},
    "label":      {"font_size": 11, "bold": True,  "color": "textLight"},
    "title_cover":{"font_size": 36, "bold": True,  "color": "white"},
    "subtitle_cover": {"font_size": 18, "bold": False, "color": "white"},
}

SLIDE_W: float = 13.33   # inches (16:9)
SLIDE_H: float = 7.5
SAFE_XMAX: float = 12.9
SAFE_YMAX: float = 7.1
CONTENT_TOP: float = 1.4   # ヘッダーバンド下端

DEFAULT_FONT = "Meiryo UI"

# ---------------------------------------------------------------------------
# 内部ユーティリティ
# ---------------------------------------------------------------------------

def _rgb(hex_str: str) -> RGBColor:
    """'RRGGBB' → RGBColor"""
    hex_str = hex_str.lstrip("#")
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def _color(key_or_hex: str) -> RGBColor:
    """カラーキー or '#RRGGBB' / 'RRGGBB' を RGBColor に変換"""
    if key_or_hex in C:
        return _rgb(C[key_or_hex])
    return _rgb(key_or_hex)


def _in(v: float) -> Emu:
    return Inches(v)


def _add_shape_fill(shape, fill_color: str | None, border_color: str | None = None, border_width: float = 1.0):
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _color(fill_color)
    else:
        shape.fill.background()

    line = shape.line
    if border_color is not None:
        line.color.rgb = _color(border_color)
        line.width = Pt(border_width)
    else:
        line.fill.background()


def _set_para_align(para, align: str):
    mapping = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    para.alignment = mapping.get(align, PP_ALIGN.LEFT)


def _apply_run_style(run, style: str | None, color: str | None, bold: bool | None, italic: bool | None, font_size: int | None):
    s = TEXT_STYLES.get(style or "body", TEXT_STYLES["body"])
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(font_size if font_size is not None else s["font_size"])
    run.font.bold = bold if bold is not None else s["bold"]
    run.font.italic = italic if italic is not None else False
    c = color or s.get("color", "text")
    run.font.color.rgb = _color(c)


# ---------------------------------------------------------------------------
# テキスト追加
# ---------------------------------------------------------------------------

def add_text(
    slide,
    x: float, y: float, w: float, h: float,
    text: str,
    style: str = "body",
    align: str = "left",
    color: str | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    font_size: int | None = None,
    word_wrap: bool = True,
) -> Any:
    """テキストボックスを追加する"""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    para = tf.paragraphs[0]
    _set_para_align(para, align)
    run = para.add_run()
    run.text = text
    _apply_run_style(run, style, color, bold, italic, font_size)
    return txBox


def add_rich_text(
    slide,
    x: float, y: float, w: float, h: float,
    paragraphs: list[dict],
    word_wrap: bool = True,
) -> Any:
    """
    複数段落・スタイル混在テキストを追加する。
    paragraphs: [{"text": str, "style": "body", "align": "left", "color": None, "bold": None, "space_before": 0}, ...]
    """
    txBox = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    first = True
    for p in paragraphs:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()

        _set_para_align(para, p.get("align", "left"))
        space_before = p.get("space_before", 0)
        if space_before:
            para.space_before = Pt(space_before)

        run = para.add_run()
        run.text = p.get("text", "")
        _apply_run_style(
            run,
            p.get("style", "body"),
            p.get("color"),
            p.get("bold"),
            p.get("italic"),
            p.get("font_size"),
        )
    return txBox


# ---------------------------------------------------------------------------
# 図形追加
# ---------------------------------------------------------------------------

def add_rect(
    slide,
    x: float, y: float, w: float, h: float,
    fill: str | None = None,
    border: str | None = None,
    border_width: float = 1.0,
    radius: float = 0.0,
) -> Any:
    """矩形を追加する。radius > 0 で角丸"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    if radius > 0:
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE = 1 in freeform, use autoshape
            _in(x), _in(y), _in(w), _in(h)
        )
    else:
        shape = slide.shapes.add_shape(
            1, _in(x), _in(y), _in(w), _in(h)
        )
    _add_shape_fill(shape, fill, border, border_width)
    return shape


def add_rounded_rect(
    slide,
    x: float, y: float, w: float, h: float,
    fill: str | None = None,
    border: str | None = None,
    border_width: float = 0.75,
) -> Any:
    """角丸矩形を追加する"""
    from pptx.util import Pt
    from lxml import etree
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
        _in(x), _in(y), _in(w), _in(h)
    )
    _add_shape_fill(shape, fill, border, border_width)
    return shape


def add_line(
    slide,
    x1: float, y1: float,
    x2: float, y2: float,
    color: str = "border",
    width: float = 1.0,
) -> Any:
    """直線を追加する"""
    from pptx.util import Pt
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        _in(x1), _in(y1), _in(x2), _in(y2)
    )
    connector.line.color.rgb = _color(color)
    connector.line.width = Pt(width)
    return connector


def add_arrow(
    slide,
    x1: float, y1: float,
    x2: float, y2: float,
    color: str = "accent",
    width: float = 2.0,
) -> Any:
    """矢印コネクタを追加する"""
    from pptx.util import Pt
    connector = slide.shapes.add_connector(
        1,
        _in(x1), _in(y1), _in(x2), _in(y2)
    )
    connector.line.color.rgb = _color(color)
    connector.line.width = Pt(width)
    # 先端矢印
    from pptx.oxml.ns import qn
    ln = connector.line._ln
    tailEnd = ln.find(qn("a:tailEnd"))
    if tailEnd is None:
        from lxml import etree
        tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "none")
    headEnd = ln.find(qn("a:headEnd"))
    if headEnd is None:
        from lxml import etree
        headEnd = etree.SubElement(ln, qn("a:headEnd"))
    headEnd.set("type", "arrow")
    headEnd.set("w", "med")
    headEnd.set("len", "med")
    return connector


def add_card(
    slide,
    x: float, y: float, w: float, h: float,
    title: str = "",
    body: str = "",
    accent: str | None = None,
    bg: str = "bg",
) -> tuple:
    """タイトル+本文カードを追加する"""
    rect = add_rounded_rect(slide, x, y, w, h, fill=bg, border="border")
    shapes = [rect]
    acc = accent or "accent"
    # アクセントバー（左端）
    bar = add_rect(slide, x, y + 0.05, 0.06, h - 0.1, fill=acc)
    shapes.append(bar)
    if title:
        t = add_text(slide, x + 0.15, y + 0.1, w - 0.2, 0.4, title, style="subheading")
        shapes.append(t)
    if body:
        b = add_text(slide, x + 0.15, y + 0.55, w - 0.2, h - 0.65, body, style="body", word_wrap=True)
        shapes.append(b)
    return tuple(shapes)


def add_pill(
    slide,
    x: float, y: float,
    text: str,
    color: str = "accent",
    text_color: str = "white",
) -> tuple:
    """ピル（小さな丸みのあるバッジ）を追加する"""
    w = max(1.2, len(text) * 0.15 + 0.4)
    h = 0.35
    rect = add_rounded_rect(slide, x, y, w, h, fill=color)
    t = add_text(slide, x + 0.05, y + 0.03, w - 0.1, h - 0.06, text,
                 style="small", align="center", color=text_color, bold=True)
    return rect, t


def add_badge(
    slide,
    x: float, y: float,
    number: int | str,
    color: str = "accent",
) -> tuple:
    """丸い番号バッジを追加する"""
    size = 0.45
    circle = slide.shapes.add_shape(9, _in(x), _in(y), _in(size), _in(size))  # OVAL
    _add_shape_fill(circle, color, None)
    t = add_text(slide, x + 0.02, y + 0.05, size - 0.04, size - 0.1,
                 str(number), style="body", align="center", color="white", bold=True, font_size=13)
    return circle, t


def add_quote(
    slide,
    x: float, y: float, w: float,
    text: str,
    source: str = "",
    bg: str = "bgLight",
) -> list:
    """引用ブロックを追加する"""
    h = 1.0 + len(text) / 80 * 0.3
    shapes = []
    rect = add_rounded_rect(slide, x, y, w, h, fill=bg)
    shapes.append(rect)
    bar = add_rect(slide, x, y, 0.08, h, fill="accent")
    shapes.append(bar)
    qt = add_text(slide, x + 0.18, y + 0.1, w - 0.28, h - 0.2,
                  f'"{text}"', style="body", italic=True)
    shapes.append(qt)
    if source:
        src = add_text(slide, x + 0.18, y + h - 0.32, w - 0.28, 0.28,
                       f"— {source}", style="small", align="right", color="textMuted")
        shapes.append(src)
    return shapes


# ---------------------------------------------------------------------------
# データ可視化（図形ベース）
# ---------------------------------------------------------------------------

def add_kpi_row(
    slide,
    y: float,
    items: list[dict],
    x_start: float = 0.5,
    total_w: float = 12.33,
) -> list:
    """
    KPI 数値カードを横並びに追加する。
    items: [{"label": str, "value": str, "unit": str, "delta": str, "color": "accent"}]
    """
    n = len(items)
    if n == 0:
        return []
    gap = 0.2
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 1.6
    shapes = []
    for i, item in enumerate(items):
        cx = x_start + i * (card_w + gap)
        rect = add_rounded_rect(slide, cx, y, card_w, card_h, fill="bg", border="border")
        shapes.append(rect)
        acc = item.get("color", "accent")
        bar = add_rect(slide, cx, y, card_w, 0.07, fill=acc)
        shapes.append(bar)
        value = item.get("value", "—")
        unit = item.get("unit", "")
        label = item.get("label", "")
        delta = item.get("delta", "")
        # 数値
        shapes.append(add_text(slide, cx + 0.1, y + 0.2, card_w - 0.2, 0.7,
                               value, style="kpi", align="center",
                               color=acc, font_size=min(40, max(20, 40 - max(0, len(value) - 4) * 4))))
        # 単位
        if unit:
            shapes.append(add_text(slide, cx + 0.1, y + 0.85, card_w - 0.2, 0.25,
                                   unit, style="small", align="center", color="textLight"))
        # ラベル
        shapes.append(add_text(slide, cx + 0.1, y + 1.1, card_w - 0.2, 0.28,
                               label, style="caption", align="center"))
        # デルタ
        if delta:
            delta_color = "accent2" if str(delta).startswith("+") else "accent3"
            shapes.append(add_text(slide, cx + 0.1, y + 1.35, card_w - 0.2, 0.22,
                                   delta, style="caption", align="center", color=delta_color, bold=True))
    return shapes


def add_table_shapes(
    slide,
    x: float, y: float, w: float,
    headers: list[str],
    rows: list[list[str]],
    col_widths: list[float] | None = None,
    row_h: float = 0.38,
    header_color: str = "headerBg",
    header_text_color: str = "headerText",
    alt_row: bool = True,
) -> list:
    """
    図形ベースのテーブルを追加する（pptx native table 禁止のため）。
    """
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols

    shapes = []
    # ヘッダー行
    cx = x
    for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
        rect = add_rect(slide, cx, y, cw, row_h, fill=header_color)
        shapes.append(rect)
        shapes.append(add_text(slide, cx + 0.08, y + 0.05, cw - 0.16, row_h - 0.1,
                               hdr, style="small", align="center",
                               color=header_text_color, bold=True))
        cx += cw

    # データ行
    for i, row in enumerate(rows):
        ry = y + row_h * (i + 1)
        bg = "rowAlt" if alt_row and i % 2 == 1 else "bg"
        cx = x
        for j, (cell, cw) in enumerate(zip(row, col_widths)):
            rect = add_rect(slide, cx, ry, cw, row_h, fill=bg, border="border")
            shapes.append(rect)
            shapes.append(add_text(slide, cx + 0.08, ry + 0.05, cw - 0.16, row_h - 0.1,
                                   str(cell), style="small", align="left"))
            cx += cw

    return shapes


def add_comparison(
    slide,
    y: float,
    left_title: str,
    right_title: str,
    left_items: list[str],
    right_items: list[str],
    x_start: float = 0.5,
    total_w: float = 12.33,
    left_color: str = "accent",
    right_color: str = "accent2",
) -> list:
    """左右2カラム比較ブロックを追加する"""
    gap = 0.3
    col_w = (total_w - gap) / 2
    shapes = []

    for idx, (title, items, color) in enumerate([
        (left_title, left_items, left_color),
        (right_title, right_items, right_color),
    ]):
        cx = x_start + idx * (col_w + gap)
        h = 0.5 + len(items) * 0.45 + 0.2
        rect = add_rounded_rect(slide, cx, y, col_w, h, fill="bgLight", border="border")
        shapes.append(rect)
        bar = add_rect(slide, cx, y, col_w, 0.08, fill=color)
        shapes.append(bar)
        shapes.append(add_text(slide, cx + 0.15, y + 0.12, col_w - 0.3, 0.38,
                               title, style="subheading", color=color))
        for k, item in enumerate(items):
            iy = y + 0.55 + k * 0.4
            shapes.append(add_text(slide, cx + 0.15, iy, col_w - 0.3, 0.38,
                                   f"• {item}", style="body"))
    return shapes


def add_content_box_table(
    slide,
    x: float, y: float, w: float,
    headers: list[str],
    rows: list[dict],
    col_widths: list[float] | None = None,
    label_w: float = 1.5,
    header_h: float = 0.55,
    row_h: float = 0.48,
    gap: float = 0.15,
    header_color: str = "accent",
    header_text_color: str = "white",
    highlight_col: int | None = None,
    highlight_label: str = "推奨",
    show_labels: bool = True,
) -> list:
    """
    コンテンツボックス型テーブルを追加する。
    各列が独立したカード状ボックスで構成される、コンサル資料でよく使われる比較表。

    headers: 列ヘッダー文字列のリスト
    rows:    [{"label": str, "values": [str, str, ...], "label_color": str}]
             label を省略するとラベル列を表示しない（show_labels=False 時も同様）
    label_w: 左ラベル列の幅（show_labels=False の場合は無視）
    highlight_col: ハイライトする列インデックス（0始まり、None でハイライトなし）
    highlight_label: ハイライト列上部に表示するバッジテキスト

    例:
        add_content_box_table(
            slide, x=0.5, y=1.6, w=12.33,
            headers=["施策A", "施策B", "施策C"],
            rows=[
                {"label": "コスト",   "values": ["高い",  "中程度", "低い"]},
                {"label": "実施期間", "values": ["6ヶ月", "3ヶ月",  "1ヶ月"]},
                {"label": "効果",     "values": ["大",    "中",     "小"]},
            ],
            highlight_col=2,
        )
    """
    n_cols = len(headers)
    n_rows = len(rows)
    if n_cols == 0:
        return []

    # 列幅計算
    content_x = x + (label_w + gap if show_labels else 0)
    content_w = w - (label_w + gap if show_labels else 0)
    if col_widths is None:
        cw = (content_w - gap * (n_cols - 1)) / n_cols
        col_widths = [cw] * n_cols

    total_h = header_h + row_h * n_rows
    shapes = []

    # ---- ハイライト列の背景（ボックス全体を浮かせる効果） ----
    if highlight_col is not None and 0 <= highlight_col < n_cols:
        hx = content_x + sum(col_widths[:highlight_col]) + gap * highlight_col
        hw = col_widths[highlight_col]
        shadow = add_rect(slide, hx + 0.04, y + 0.04, hw, total_h, fill="border")
        shapes.append(shadow)
        # バッジ
        badge_w = min(hw * 0.7, 1.5)
        badge_x = hx + (hw - badge_w) / 2
        badge = add_rounded_rect(slide, badge_x, y - 0.3, badge_w, 0.28, fill=header_color)
        shapes.append(badge)
        shapes.append(add_text(slide, badge_x + 0.05, y - 0.28, badge_w - 0.1, 0.24,
                               highlight_label, style="caption", align="center",
                               color="white", bold=True))

    # ---- 左ラベル列 ----
    if show_labels:
        # ラベル列ヘッダー（空白）
        shapes.append(add_rounded_rect(slide, x, y, label_w, header_h, fill="bgLight", border="border"))
        # 各行ラベル
        for i, row in enumerate(rows):
            ry = y + header_h + row_h * i
            lbl = row.get("label", "")
            lbl_color = row.get("label_color", "textLight")
            shapes.append(add_rect(slide, x, ry, label_w, row_h, fill="bgLight", border="border"))
            if lbl:
                shapes.append(add_text(slide, x + 0.12, ry + 0.06, label_w - 0.2, row_h - 0.1,
                                       lbl, style="small", align="left",
                                       bold=True, color=lbl_color))

    # ---- 各データ列 ----
    for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
        cx = content_x + sum(col_widths[:j]) + gap * j
        is_highlight = highlight_col is not None and j == highlight_col
        h_color = header_color if not is_highlight else header_color
        box_bg = "bg" if not is_highlight else "bg"
        h_border = header_color if is_highlight else "border"

        # ヘッダーボックス
        shapes.append(add_rounded_rect(
            slide, cx, y, cw, header_h,
            fill=h_color if is_highlight else "bgLight",
            border=h_color if is_highlight else "border",
        ))
        shapes.append(add_text(
            slide, cx + 0.1, y + 0.08, cw - 0.2, header_h - 0.14,
            hdr, style="small", align="center",
            color=header_text_color if is_highlight else "text",
            bold=True,
        ))

        # データ行
        for i, row in enumerate(rows):
            ry = y + header_h + row_h * i
            val = row.get("values", [""] * n_cols)
            val_text = val[j] if j < len(val) else ""

            # セル背景
            cell_fill = "rowAlt" if (not is_highlight and i % 2 == 1) else "bg"
            if is_highlight:
                cell_fill = "bg"
            shapes.append(add_rect(
                slide, cx, ry, cw, row_h,
                fill=cell_fill,
                border=header_color if is_highlight else "border",
            ))

            # セルテキスト
            shapes.append(add_text(
                slide, cx + 0.1, ry + 0.06, cw - 0.2, row_h - 0.1,
                val_text, style="small", align="center",
                color="text" if not is_highlight else "text",
            ))

    return shapes


def add_gantt(
    slide,
    x: float, y: float, w: float,
    periods: list[dict],
    rows: list[dict],
    legend: dict | None = None,
    header_color: str = "accent",
    row_h: float = 0.40,
    label_w: float = 1.0,
    group_w: float = 0.45,
    bar_colors: dict | None = None,
) -> list:
    """
    コンサル式ガントチャート（プロジェクトスケジュール）を追加する。

    periods: 期間ヘッダー定義
        [{"label": "10月", "subperiods": ["3週", "10週", "17週", "24週", "31週"]}, ...]
        サブ期間なしの場合: [{"label": "Phase1", "cols": 3}, ...]

    rows: 行定義
        [
          # マイルストン行
          {"label": "マイルストン", "milestones": [
              {"col": 0, "label": "契約開始"},
              {"col": 5, "label": "キックオフ"},
          ]},
          # タスク行（グループあり）
          {"label": "機能調査\\n拡張", "group": "調査研究", "tasks": [
              {"start": 2, "end": 5,  "label": "調査対象機能の抽出", "color": "self"},
              {"start": 6, "end": 10, "label": "機能の選定",          "color": "self"},
          ]},
          # グループなしタスク行
          {"label": "資料作成", "tasks": [
              {"start": 18, "end": 20, "label": "報告書作成", "color": "joint"},
          ]},
        ]

    legend: {"弊社作業": "self", "貴社作業": "client", "合同": "joint"}
    bar_colors: {"self": "accent", "client": "accent3", "joint": "accent2"}
    """
    if bar_colors is None:
        bar_colors = {
            "self":   C["accent"],
            "client": C["accent3"],
            "joint":  C["accent2"],
        }

    # 総カラム数
    def _period_cols(p):
        return len(p.get("subperiods", [])) or p.get("cols", 1)
    total_cols = sum(_period_cols(p) for p in periods)

    # グループ列の有無
    has_groups = any(r.get("group") for r in rows)
    gw = group_w if has_groups else 0.0

    header_h1 = 0.30   # 月ヘッダー高
    header_h2 = 0.24   # 週サブヘッダー高
    has_sub = any(p.get("subperiods") for p in periods)
    header_h = header_h1 + (header_h2 if has_sub else 0)

    content_x = x + label_w + gw
    content_w = w - label_w - gw
    col_w = content_w / total_cols

    shapes = []

    # ---- 空白の左上コーナー ----
    shapes.append(add_rect(slide, x, y, label_w + gw, header_h, fill="bgLight", border="border"))

    # ---- 月ヘッダー ----
    cx = content_x
    for p in periods:
        n = _period_cols(p)
        pw = col_w * n
        shapes.append(add_rect(slide, cx, y, pw, header_h1, fill=header_color))
        shapes.append(add_text(slide, cx + 0.04, y + 0.04, pw - 0.06, header_h1 - 0.06,
                               p["label"], style="caption", align="center",
                               color="white", bold=True))
        cx += pw

    # ---- サブ期間（週）ヘッダー ----
    if has_sub:
        cx = content_x
        for p in periods:
            for sub in p.get("subperiods", []):
                shapes.append(add_rect(slide, cx, y + header_h1, col_w, header_h2,
                                       fill="bgLight", border="border"))
                shapes.append(add_text(slide, cx + 0.01, y + header_h1 + 0.02,
                                       col_w - 0.02, header_h2 - 0.04,
                                       sub, style="caption", align="center",
                                       color="textMuted", font_size=9))
                cx += col_w

    # ---- グループラベル（縦結合） ----
    if has_groups:
        i = 0
        while i < len(rows):
            g = rows[i].get("group")
            if g:
                j = i + 1
                while j < len(rows) and rows[j].get("group") == g:
                    j += 1
                gy = y + header_h + i * row_h
                gh = (j - i) * row_h
                shapes.append(add_rect(slide, x + label_w, gy, gw, gh,
                                       fill=header_color, border=None))
                # テキスト：中央に縦書き風（短い場合）/ 横書き小フォント
                shapes.append(add_text(slide, x + label_w + 0.02, gy + gh / 2 - 0.25,
                                       gw - 0.04, 0.5, g,
                                       style="caption", align="center",
                                       color="white", bold=True, font_size=9))
                i = j
            else:
                i += 1

    # ---- 行ラベル + グリッド + バー + マイルストン ----
    for i, row in enumerate(rows):
        ry = y + header_h + i * row_h

        # 行ラベル
        shapes.append(add_rect(slide, x, ry, label_w, row_h,
                               fill="bgLight", border="border"))
        shapes.append(add_text(slide, x + 0.07, ry + 0.04, label_w - 0.1, row_h - 0.06,
                               row["label"], style="caption", align="left",
                               font_size=9, bold=False))

        # グリッド背景
        row_bg = "bg" if i % 2 == 0 else "rowAlt"
        shapes.append(add_rect(slide, content_x, ry, content_w, row_h,
                               fill=row_bg, border="border"))

        # 縦グリッド線
        for ci in range(1, total_cols):
            lx = content_x + ci * col_w
            shapes.append(add_line(slide, lx, ry, lx, ry + row_h,
                                   color="border", width=0.3))

        # タスクバー
        for task in row.get("tasks", []):
            s = max(0, task["start"])
            e = min(total_cols - 1, task["end"])  # 範囲外クランプ
            bar_x = content_x + s * col_w + 0.03
            bar_w = (e - s + 1) * col_w - 0.06
            bar_y = ry + row_h * 0.18
            bar_h = row_h * 0.64
            tc = task.get("color", "self")
            fill = bar_colors.get(tc, tc)
            shapes.append(add_rect(slide, bar_x, bar_y, bar_w, bar_h, fill=fill))
            if task.get("label") and bar_w > 0.3:
                shapes.append(add_text(slide, bar_x + 0.05, bar_y + 0.02,
                                       bar_w - 0.08, bar_h - 0.04,
                                       task["label"], style="caption",
                                       align="left", color="white", font_size=9))

        # マイルストン（▼マーカー）
        for ms in row.get("milestones", []):
            col = max(0, min(total_cols - 1, ms["col"]))  # 範囲外クランプ
            mx = content_x + (col + 0.5) * col_w
            shapes.append(add_text(slide, mx - 0.18, ry + 0.01, 0.36, row_h * 0.45,
                                   "▼", style="caption", align="center",
                                   color=header_color, bold=True, font_size=10))
            if ms.get("label"):
                lbl_w = max(1.0, len(ms["label"]) * 0.13)
                # 右端クリップ
                lbl_x = max(x, min(mx - lbl_w / 2, x + w - lbl_w))
                shapes.append(add_text(slide, lbl_x, ry + row_h * 0.42,
                                       lbl_w, row_h * 0.55,
                                       ms["label"], style="caption",
                                       align="center", color="text",
                                       font_size=9))

    # ---- 凡例 ----
    if legend:
        n_items = len(legend)
        lw = 1.7
        lh = 0.22 * n_items + 0.18
        legend_x = content_x + content_w - lw - 0.1
        legend_y = y + header_h + 0.08
        shapes.append(add_rounded_rect(slide, legend_x, legend_y, lw, lh,
                                       fill="bg", border="border"))
        shapes.append(add_text(slide, legend_x + 0.1, legend_y + 0.03, lw - 0.2, 0.18,
                               "凡例：", style="caption", color="textMuted",
                               font_size=9, bold=True))
        for li, (lbl, color_key) in enumerate(legend.items()):
            ly = legend_y + 0.18 + li * 0.22
            fill = bar_colors.get(color_key, color_key)
            shapes.append(add_rect(slide, legend_x + 0.1, ly + 0.03, 0.22, 0.14, fill=fill))
            shapes.append(add_text(slide, legend_x + 0.38, ly, lw - 0.48, 0.22,
                                   lbl, style="caption", color="text", font_size=9))

    return shapes


def add_bar_chart_img(
    slide,
    x: float, y: float, w: float, h: float,
    data: dict,
    title: str = "",
) -> Any:
    """
    matplotlib で横棒グラフを生成し PNG として埋め込む。
    data: {"labels": [...], "values": [...], "colors": [...] (optional)}
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm

    # 日本語フォント
    font_paths = [
        "/mnt/c/Windows/Fonts/meiryo.ttc",
        "/mnt/c/Windows/Fonts/YuGothM.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    ]
    for fp in font_paths:
        if os.path.exists(fp):
            prop = fm.FontProperties(fname=fp)
            plt.rcParams["font.family"] = prop.get_name()
            break

    labels = data.get("labels", [])
    values = data.get("values", [])
    colors_list = data.get("colors", [C["accent"]] * len(labels))
    colors_rgb = [f"#{c.lstrip('#')}" for c in colors_list]

    fig_w = w * 1.5
    fig_h = max(2.0, len(labels) * 0.4 + 0.8)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.barh(labels, values, color=colors_rgb, edgecolor="none")
    ax.set_xlabel("")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.tick_params(left=False)
    if title:
        ax.set_title(title, fontsize=11)
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, _in(x), _in(y), _in(w), _in(h))
    return pic


# ---------------------------------------------------------------------------
# フロー・構造
# ---------------------------------------------------------------------------

def add_flow_row(
    slide,
    y: float,
    steps: list[str],
    color: str = "accent",
    x_start: float = 0.5,
    total_w: float = 12.33,
    box_h: float = 0.9,
) -> list:
    """横並びフローステップ（箱＋矢印）を追加する"""
    n = len(steps)
    if n == 0:
        return []
    arrow_w = 0.3
    box_w = (total_w - arrow_w * (n - 1)) / n
    shapes = []

    for i, step in enumerate(steps):
        cx = x_start + i * (box_w + arrow_w)
        rect = add_rounded_rect(slide, cx, y, box_w, box_h, fill=color)
        shapes.append(rect)
        shapes.append(add_text(slide, cx + 0.05, y + 0.05, box_w - 0.1, box_h - 0.1,
                               step, style="body", align="center", color="white", bold=True))
        if i < n - 1:
            ax_start = cx + box_w
            ay = y + box_h / 2
            shapes.append(add_arrow(slide, ax_start, ay, ax_start + arrow_w, ay, color=color))
    return shapes


def add_timeline(
    slide,
    y: float,
    events: list[dict],
    x_start: float = 0.5,
    total_w: float = 12.33,
    color: str = "accent",
) -> list:
    """
    横タイムラインを追加する。
    events: [{"date": str, "title": str, "desc": str}]
    """
    n = len(events)
    if n == 0:
        return []
    shapes = []
    line_y = y + 0.5
    shapes.append(add_line(slide, x_start, line_y, x_start + total_w, line_y, color=color, width=2))

    spacing = total_w / n
    for i, ev in enumerate(events):
        cx = x_start + i * spacing + spacing / 2
        # 点
        dot = slide.shapes.add_shape(9, _in(cx - 0.12), _in(line_y - 0.12), _in(0.24), _in(0.24))
        _add_shape_fill(dot, color)
        shapes.append(dot)
        # 日付（上）
        shapes.append(add_text(slide, cx - spacing * 0.45, line_y - 0.5, spacing * 0.9, 0.35,
                               ev.get("date", ""), style="caption", align="center", color=color, bold=True))
        # タイトル（下）
        shapes.append(add_text(slide, cx - spacing * 0.45, line_y + 0.18, spacing * 0.9, 0.4,
                               ev.get("title", ""), style="small", align="center", bold=True))
        # 説明
        if ev.get("desc"):
            shapes.append(add_text(slide, cx - spacing * 0.45, line_y + 0.58, spacing * 0.9, 0.6,
                                   ev.get("desc", ""), style="caption", align="center"))
    return shapes


def add_agenda(
    slide,
    items: list[str],
    current: int | None = None,
    x_start: float = 2.0,
    y_start: float = 1.8,
    item_h: float = 0.65,
    w: float = 9.0,
    color: str = "accent",
) -> list:
    """アジェンダリストを追加する。current はハイライトするインデックス（0始まり）"""
    shapes = []
    for i, item in enumerate(items):
        is_current = current is not None and i == current
        iy = y_start + i * item_h
        bg = color if is_current else "bgLight"
        tc = "white" if is_current else "text"
        rect = add_rounded_rect(slide, x_start, iy, w, item_h - 0.1, fill=bg)
        shapes.append(rect)
        num_bg = "white" if is_current else color
        num_tc = color if is_current else "white"
        # 番号
        c, t = add_badge(slide, x_start + 0.18, iy + 0.08, i + 1, color=num_bg)
        shapes.extend([c, t])
        # テキスト
        shapes.append(add_text(slide, x_start + 0.8, iy + 0.1, w - 1.0, item_h - 0.2,
                               item, style="body", color=tc, bold=is_current))
    return shapes


# ---------------------------------------------------------------------------
# ヘッダー・フッター
# ---------------------------------------------------------------------------

def add_slide_header(
    slide,
    title: str,
    section: str = "",
    color: str = "accent",
) -> list:
    """スライドヘッダーバンドを追加する"""
    shapes = []
    bar = add_rect(slide, 0, 0, SLIDE_W, 1.1, fill=color)
    shapes.append(bar)
    if section:
        shapes.append(add_text(slide, 0.4, 0.1, SLIDE_W - 0.8, 0.35,
                               section, style="small", color="white", bold=False))
    shapes.append(add_text(slide, 0.4, 0.38, SLIDE_W - 0.8, 0.65,
                           title, style="heading", color="white", font_size=22))
    return shapes


def add_slide_footer(
    slide,
    page_num: int | None = None,
    footer_text: str = "",
) -> list:
    """スライドフッターを追加する"""
    shapes = []
    shapes.append(add_line(slide, 0.4, SLIDE_H - 0.4, SLIDE_W - 0.4, SLIDE_H - 0.4,
                           color="border", width=0.5))
    if footer_text:
        shapes.append(add_text(slide, 0.4, SLIDE_H - 0.38, 8.0, 0.3,
                               footer_text, style="caption", color="textMuted"))
    if page_num is not None:
        shapes.append(add_text(slide, SLIDE_W - 1.4, SLIDE_H - 0.38, 1.0, 0.3,
                               str(page_num), style="caption", align="right", color="textMuted"))
    return shapes


# ---------------------------------------------------------------------------
# SlideBuilder
# ---------------------------------------------------------------------------

class SlideBuilder:
    """
    PPTXプレゼンテーションをビルドするクラス。
    theme: {"accent": "RRGGBB", "accent2": "RRGGBB", ...} でカラーを上書き可能。
    """

    def __init__(
        self,
        theme: dict | str | None = None,
        master_path: str | None = None,
    ):
        # テーマ上書き（毎回デフォルトから作り直してインスタンス間で色が混ざらないようにする）
        global C
        if isinstance(theme, str):
            theme = THEMES.get(theme, {})
        C = {**_C_DEFAULT, **(theme or {})}

        # プレゼンテーション初期化
        if master_path and Path(master_path).exists():
            self.prs = Presentation(master_path)
        else:
            self.prs = Presentation()

        # スライドサイズ設定
        self.prs.slide_width = _in(SLIDE_W)
        self.prs.slide_height = _in(SLIDE_H)

        self._page = 0
        self._footer_text = ""

    def set_footer(self, text: str):
        self._footer_text = text

    def _new_blank_slide(self):
        """空白レイアウトでスライドを追加する"""
        layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(layout)
        # 既存プレースホルダを削除
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        self._page += 1
        return slide

    def add_cover(
        self,
        title: str,
        subtitle: str = "",
        date: str = "",
        bg_color: str | None = None,
    ):
        """表紙スライドを追加する"""
        slide = self._new_blank_slide()
        bg = bg_color or C["accent"]
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=bg)
        # タイトルボックス
        add_rect(slide, 0, SLIDE_H / 2 - 1.6, SLIDE_W, 3.2, fill=None)
        add_text(slide, 1.0, SLIDE_H / 2 - 1.4, SLIDE_W - 2.0, 1.8,
                 title, style="title_cover", align="center", font_size=min(40, max(24, 48 - len(title) // 3)))
        if subtitle:
            add_text(slide, 1.0, SLIDE_H / 2 + 0.5, SLIDE_W - 2.0, 0.6,
                     subtitle, style="subtitle_cover", align="center")
        if date:
            add_text(slide, 1.0, SLIDE_H - 1.0, SLIDE_W - 2.0, 0.4,
                     date, style="caption", align="right", color="white")
        # デコレーション下バー
        add_rect(slide, 0, SLIDE_H - 0.25, SLIDE_W, 0.25, fill="white")
        return slide

    def add_section(
        self,
        title: str,
        subtitle: str = "",
        number: int | str | None = None,
        color: str | None = None,
    ):
        """セクション区切りスライドを追加する"""
        slide = self._new_blank_slide()
        acc = color or C["accent"]
        # 左帯
        add_rect(slide, 0, 0, 4.5, SLIDE_H, fill=acc)
        add_rect(slide, 4.5, 0, SLIDE_W - 4.5, SLIDE_H, fill="bgLight")
        if number is not None:
            add_text(slide, 0.4, SLIDE_H / 2 - 1.2, 3.7, 1.0,
                     str(number), style="kpi", align="center", color="white", font_size=72)
        add_text(slide, 0.4, SLIDE_H / 2 - 0.1, 3.7, 0.8,
                 "SECTION", style="caption", align="center", color="white")
        add_text(slide, 5.0, SLIDE_H / 2 - 0.8, SLIDE_W - 5.5, 1.0,
                 title, style="heading", align="left", color=acc, font_size=28)
        if subtitle:
            add_text(slide, 5.0, SLIDE_H / 2 + 0.3, SLIDE_W - 5.5, 0.6,
                     subtitle, style="body", align="left", color="textLight")
        add_slide_footer(slide, self._page, self._footer_text)
        return slide

    def add_body(
        self,
        title: str,
        section: str = "",
        color: str | None = None,
    ):
        """本文スライドのベース（ヘッダー付き）を追加する"""
        slide = self._new_blank_slide()
        acc = color or C["accent"]
        add_slide_header(slide, title, section, color=acc)
        add_slide_footer(slide, self._page, self._footer_text)
        return slide

    def save(self, path: str) -> None:
        """PPTXを保存する"""
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        self.prs.save(path)

    def save_and_validate(self, path: str) -> "ValidationResult":
        """PPTXを保存してバリデーションを実行する"""
        self.save(path)
        result = validate(path)
        if result.errors:
            print(f"[VALIDATION ERRORS] {path}")
            for e in result.errors:
                print(f"  ERROR: {e}")
        if result.warnings:
            for w in result.warnings:
                print(f"  WARN:  {w}")
        return result


# ---------------------------------------------------------------------------
# バリデーション
# ---------------------------------------------------------------------------

@dataclass
class ValidationResult:
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return len(self.errors) == 0

    def __repr__(self):
        status = "OK" if self.ok else "FAIL"
        return (f"ValidationResult({status}, errors={self.errors}, warnings={self.warnings})")


def validate(pptx_path: str) -> ValidationResult:
    """PPTXファイルをバリデーションする"""
    result = ValidationResult()

    if not Path(pptx_path).exists():
        result.errors.append(f"ファイルが見つかりません: {pptx_path}")
        return result

    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        slide_no = slide_idx + 1
        shape_count = 0

        for shape in slide.shapes:
            shape_count += 1

            # pptx native table 禁止
            if shape.has_table:
                result.errors.append(f"スライド{slide_no}: pptx native table が使用されています（add_table_shapes を使用してください）")

            # テキストフレームのチェック
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        fn = run.font.name
                        if fn and fn in FORBIDDEN_FONTS:
                            result.errors.append(f"スライド{slide_no}: 禁止フォント '{fn}' が使用されています")
                        fs = run.font.size
                        if fs is not None:
                            fs_pt = fs.pt
                            if fs_pt < 9:
                                result.warnings.append(f"スライド{slide_no}: フォントサイズが小さすぎます ({fs_pt:.0f}pt < 9pt)")
                            elif fs_pt > 60:
                                result.warnings.append(f"スライド{slide_no}: フォントサイズが大きすぎます ({fs_pt:.0f}pt > 60pt)")

            # はみ出しチェック
            # 背景・ヘッダーバー・フッターライン等はチェック対象外
            try:
                left_emu = shape.left
                top_emu = shape.top
                right_edge = shape.left + shape.width
                bottom_edge = shape.top + shape.height
                left_in = left_emu / 914400 if left_emu else 0
                top_in = top_emu / 914400 if top_emu else 0

                # フルブリード・ヘッダー・フッター要素はチェック対象外
                width_in = shape.width / 914400 if shape.width else 0
                height_in = shape.height / 914400 if shape.height else 0
                # 左端から始まる or スライド幅に近い幅 → 右端チェックをスキップ
                skip_right = left_in <= 0.05 or width_in >= SLIDE_W - 0.5
                # 上端から始まる or スライド高に近い高さ or スライド下端に到達 → 下端チェックをスキップ
                reaches_bottom = (top_in + height_in) >= SLIDE_H - 0.15
                skip_bottom = top_in <= 0.05 or height_in >= SLIDE_H - 0.5 or reaches_bottom
                # ヘッダー帯 (y < CONTENT_TOP) はすべてスキップ
                is_header_zone = top_in < CONTENT_TOP
                # フッター帯 (y > SLIDE_H - 0.55) はすべてスキップ
                is_footer_zone = top_in >= SLIDE_H - 0.55

                if not is_header_zone and not is_footer_zone:
                    if not skip_right and right_edge is not None and right_edge > _in(SAFE_XMAX):
                        result.errors.append(
                            f"スライド{slide_no} '{shape.name}': 右端はみ出し "
                            f"({right_edge / 914400:.2f}\" > {SAFE_XMAX}\")"
                        )
                    if not skip_bottom and bottom_edge is not None and bottom_edge > _in(SAFE_YMAX):
                        result.errors.append(
                            f"スライド{slide_no} '{shape.name}': 下端はみ出し "
                            f"({bottom_edge / 914400:.2f}\" > {SAFE_YMAX}\")"
                        )
            except Exception:
                pass

        # 1スライド文字数
        all_text = " ".join(
            run.text
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        )
        if len(all_text) > 300:
            result.warnings.append(f"スライド{slide_no}: テキスト量が多い ({len(all_text)}文字 > 300文字)")

        # シェイプ数
        if shape_count > 25:
            result.warnings.append(f"スライド{slide_no}: シェイプ数が多い ({shape_count} > 25)")

    return result
